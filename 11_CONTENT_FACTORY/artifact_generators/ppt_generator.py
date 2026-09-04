# 11_CONTENT_FACTORY/artifact_generators/ppt_generator.py — 真实 PPTX 生成

from __future__ import annotations

from pathlib import Path
from typing import Any


def generate_ppt(title: str, sections: list[dict], style: str, output_path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        return {"status": "error", "file_path": "", "error": f"python-pptx not installed: {exc}"}

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 首页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    sub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
    sub.text_frame.text = f"Style: {style} | AI Factory OS Content Factory"

    # 目录页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(0.8))
    box.text_frame.text = "目录"
    box.text_frame.paragraphs[0].font.size = Pt(32)
    box.text_frame.paragraphs[0].font.bold = True
    y = 1.8
    for i, sec in enumerate(sections, 1):
        tb = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(10), Inches(0.5))
        tb.text_frame.text = f"{i}. {sec.get('title', 'Section')}"
        tb.text_frame.paragraphs[0].font.size = Pt(20)
        y += 0.55

    # 内容页
    for sec in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
        box.text_frame.text = sec.get("title", "内容")
        box.text_frame.paragraphs[0].font.size = Pt(28)
        box.text_frame.paragraphs[0].font.bold = True
        y = 1.6
        for bullet in sec.get("bullets", []):
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(10.5), Inches(0.45))
            tb.text_frame.text = f"• {bullet}"
            tb.text_frame.paragraphs[0].font.size = Pt(18)
            y += 0.5

    # 结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
    box.text_frame.text = "谢谢"
    box.text_frame.paragraphs[0].font.size = Pt(44)
    box.text_frame.paragraphs[0].font.bold = True

    prs.save(str(output_path))
    return {"status": "ok", "file_path": str(output_path)}
